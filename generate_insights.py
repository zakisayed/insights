import streamlit as st
import os
import json
from moviepy.editor import VideoFileClip
import assemblyai as aai
import google.generativeai as genai
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from io import BytesIO
import json
from docx import Document
import tempfile
import time
# ------------------ Utility Functions ------------------

# Load API keys from user inputs (Session State)
def load_api_keys():
    return st.session_state.get('assemblyai_key', None), st.session_state.get('gemini_key', None)

# Save API keys into session state
def save_api_keys(assemblyai_key, gemini_key):
    st.session_state['assemblyai_key'] = assemblyai_key
    st.session_state['gemini_key'] = gemini_key

# Extract audio from uploaded video
def extract_audio_from_video(uploaded_file):
    # Create a temporary file to save the uploaded video
    with tempfile.NamedTemporaryFile(delete=False, suffix=".mp4") as temp_video_file:
        temp_video_file.write(uploaded_file.read())  # Write the video content to the temp file
        temp_video_file_path = temp_video_file.name  # Get the file path of the temporary file

    # Now, use moviepy to process the video file from the temporary file
    video_clip = VideoFileClip(temp_video_file_path)
    audio_clip = video_clip.audio
    audio_file = "output_audio.mp3"
    audio_clip.write_audiofile(audio_file)
    video_clip.close()
    audio_clip.close()
    
    return audio_file  # Return the path of the audio file

# Ask Google Gemini for insights, doc, or presentation generation
def ask_gemini(api_key, transcript, prompt):
    genai.configure(api_key=api_key)  # Configure Gemini API key
    model = genai.GenerativeModel("gemini-1.5-flash")  # Set the model
    message = f"{prompt}\nTranscript: {transcript}"  # Format the prompt with transcript
    response = model.generate_content(message)
    
    try:
        response_text = response.candidates[0].content.parts[0].text.strip()  # Extract response
        return response_text
    except (AttributeError, IndexError):
        return None  # Handle errors

def generate_insights_prompt():
    return """Please summarize the following meeting in a clear, concise, and professional manner, ensuring that all key points are accurately reflected. The summary should be proofread, grammatically correct, and structured as a formal meeting summary or minutes of the meeting. The following elements must be included:

Meeting Details (if applicable otherwise populate Not Available): Include the meeting date, time, location (or virtual platform), and participants' names and roles.
Meeting Objective: Clearly state the purpose or agenda of the meeting.
Key Discussion Points: Highlight the most important topics discussed during the meeting. Include decisions made, different perspectives or insights shared by participants, and any challenges or concerns raised.
Action Items: List out all actionable tasks or next steps, specifying who is responsible for each task and the expected deadline or timeline.
Key Takeaways: Provide a brief conclusion summarizing the most critical outcomes or decisions of the meeting.
Follow-Up Meeting (if applicable): Mention if a follow-up meeting is required and note any relevant dates or action items to be revisited.
Tone and Language: Ensure the tone is formal and professional. The language should be polished, free from errors, and avoid any jargon unless specific to the context of the meeting.

For example:

Meeting Details: Date: [Insert Date], Time: [Insert Time], Participants: [Names and Roles].
Objective: The primary goal of this meeting was to discuss [specific agenda or issue].
Discussion Points: [Insert details about major topics].
Decisions Made: [Summarize key decisions].
Action Items: [List of tasks, responsible persons, and deadlines].
Key Takeaways: [Brief wrap-up with important conclusions].
Next Steps/Follow-Up: [Details about any upcoming meetings or tasks]."""

def generate_doc_prompt():
    return """You are provided with the transcript of a meeting. Your task is to generate a structured and professional Requirement Document based on the discussion in the transcript. The document should clearly outline the objectives, scope, requirements, and other relevant details necessary for the successful completion of the project. Use the following structure and guidelines:

1. Title Page:
Document Title: "Project Requirements Document"
Project Name: The name of the project discussed in the meeting.
Version Number: Version 1.0 (or relevant version).
Date: The date the meeting was held.
Prepared By: The name(s) of key participants or teams who contributed to the document.
2. Table of Contents:
Generate an automatic table of contents based on the document structure, including all sections and subsections with page numbers for easy navigation.
3. Executive Summary:
Provide a concise summary of the meeting, stating the purpose, goals, and key outcomes.
This section should briefly highlight the project’s primary objectives and high-level requirements.
4. Project Overview:
Introduction: A brief introduction describing the project, its background, and the overall business problem it seeks to solve.
Objectives: Clearly state the goals and outcomes expected from this project. Ensure these are specific, measurable, achievable, relevant, and time-bound (SMART).
Stakeholders: List all key stakeholders identified in the meeting, including their roles (e.g., Product Owner, Development Team, End Users, etc.).
5. Scope of Work:
In-Scope: List all the features, functionalities, and deliverables that are within the scope of the project, as discussed in the meeting.
Out-of-Scope: Highlight any areas or features that were explicitly mentioned as out-of-scope, ensuring there is no ambiguity.
6. Functional Requirements:
For each requirement, ensure the following structure:

Requirement ID: Assign a unique identifier to each functional requirement (e.g., FR-001).
Requirement Description: A clear and concise description of what the system should do, based on the discussion in the meeting.
Priority: Categorize the priority (e.g., High, Medium, Low).
Acceptance Criteria: Specify the conditions that must be met for the requirement to be accepted.
Dependencies: List any dependencies between this requirement and other features, modules, or systems.
Owner: Specify the individual or team responsible for ensuring this requirement is met.
7. Non-Functional Requirements:
Provide a detailed list of non-functional requirements discussed in the meeting, such as:

Performance Requirements: Expected system response times, load capacity, etc.
Security Requirements: Data protection, encryption, authentication, and authorization standards.
Usability Requirements: Accessibility, ease of use, interface standards.
Reliability Requirements: Uptime, recovery, and failover expectations.
Scalability Requirements: Future scalability considerations.
Compliance Requirements: Any regulatory or compliance needs that the system must adhere to.
8. Technical Requirements:
Include any specific technical constraints or requirements mentioned in the meeting, such as:

Platform Requirements: Operating systems, browsers, devices, etc.
Integration Requirements: Specify any third-party systems or APIs the project needs to integrate with.
Database Requirements: Any database architecture or storage needs.
Technology Stack: Identify any preferred programming languages, frameworks, or tools discussed.
9. User Stories (if applicable):
Based on the discussion, create user stories where relevant:

Story ID: Assign a unique identifier (e.g., US-001).
As a [user role], I want to [action], so that [desired outcome].
Acceptance Criteria: Define clear conditions for each user story that must be fulfilled for completion.
10. Assumptions and Constraints:
Assumptions: List any assumptions made during the meeting that impact the requirements, such as availability of resources, timelines, or external dependencies.
Constraints: Identify any limitations, such as budgetary, technological, or time-related constraints.
11. Risks:
Document any potential risks identified during the meeting that could impact the project’s success.
Provide a brief description of each risk along with mitigation strategies or contingency plans discussed.
12. Timeline & Milestones:
Outline key milestones, deliverable dates, or phases based on the meeting’s discussion.
Include any deadlines agreed upon during the meeting for specific features or deliverables.
13. Responsibilities:
Team Members: List all key team members and their roles, as well as the areas or deliverables they are responsible for.
Accountability: Ensure each task or deliverable has a clearly assigned owner.
14. Approval and Sign-Off:
Add a section for approval, where relevant stakeholders or decision-makers will sign off on the final requirements.
Include space for the date and name of each approver.
15. Appendices (if needed):
Add any additional information, diagrams, flowcharts, or references that were discussed in the meeting or that would support the document (e.g., business process flows, data models, etc.).
Formatting Guidelines:
Use professional and consistent formatting throughout the document.
Ensure headings are clear and content is organized in bullet points or short paragraphs for readability.
Where applicable, use tables to summarize information (e.g., requirements, dependencies, risks, responsibilities).
The goal is to provide a comprehensive, organized, and actionable requirement document that reflects all the key points from the meeting and clearly communicates what is required for the project’s success.

**Response Format**: 
Make sure to generate the response in this format:
Make sure to not include "```json" in the response. The response should purely be in the following structure without any addition or subtraction.
Remove '''json from the response. follow only the format below

# Example response from the LLM
{
    "title": "Title of the document",
    "author": "Author Name",
    "sections": [
        {
            "heading": "",
            "body": ""
        }
    ]
}
""" 

def generate_presentation_prompt():
    return """You are provided with the transcript of a meeting. Your task is to generate content for a PowerPoint presentation based on the information discussed. The presentation should be structured, professional, and engaging, following these guidelines:

1. **Title Slide**:
    - Slide Title: Include a title that represents the main topic of the meeting.
    - Subtitle: Add the date and participants of the meeting.
    - Additional Info: Suggest a relevant image or icon that visually represents the topic in a clean and professional manner.

2. **Agenda Slide**:
    - Summarize the key topics or objectives discussed in the meeting.
    - Include 4–5 bullet points that represent each major agenda item.

3. **Key Takeaways Slide**:
    - Provide a concise summary of the most important points covered in the meeting.
    - Organize these into bullet points or numbered lists, focusing on actionable insights or decisions made.

4. **Detailed Discussion/Topics Slides**:
    - For each major point or agenda item discussed, create individual slides with the following structure:
        - Slide Title: Clearly state the main topic or discussion point.
        - Subheadings/Bullet Points: Present key points discussed in short, to-the-point statements.
        - Visual Elements: If applicable, describe simple charts, graphs, or tables to present data or statistics that enhance understanding (e.g., decisions on targets, KPIs, timelines).
        - Action Items/Decisions: Mention any specific actions or decisions made during the discussion.

5. **Challenges/Concerns Slide**:
    - Highlight challenges, risks, or concerns raised during the meeting.
    - Structure the content in bullet points with brief descriptions and any possible solutions discussed.

6. **Next Steps/Action Items Slide**:
    - Include a list of follow-up actions based on the meeting, structured clearly:
        - Action Item
        - Responsible Person
        - Deadline/Timeline

7. **Conclusion/Closing Remarks Slide**:
    - Summarize the overall outcome of the meeting in a few sentences.
    - Mention major decisions or directions agreed upon and the general sentiment (e.g., positive outlook, consensus).

8. **Thank You/Contact Information Slide**:
    - End with a professional “Thank You” slide, including relevant contact information for follow-up.

**Response Format**: 
Make sure to generate the response exactly in the following template without including any images
Do not include "```json" in the response. The response should purely be in the following structure without any addition or subtraction.
Do not add any \n

# Example response from the LLM
    [
    {
        "title": "Title",
        "text": "Content"
    },
    {
        "title": "Title 2",
        "text": "Content"
    }
    ]
"""

# Handle transcription using AssemblyAI
def transcribe_audio(assemblyai_key, audio_file):
    aai.settings.api_key = assemblyai_key  # Set AssemblyAI API key
    transcriber = aai.Transcriber()  # Create transcriber
    speaker_config = aai.TranscriptionConfig(speaker_labels=True)  # Enable speaker labeling
    transcript_data = transcriber.transcribe(audio_file, config=speaker_config)  # Transcribe audio
    
    # Format transcript with speaker labels
    transcript = "\n".join([f"Speaker {utt.speaker}: {utt.text}" for utt in transcript_data.utterances])
    return transcript

# Handle file processing (video/text/docx) and generating insights, docs, or presentations
def process_file(file, assemblyai_key, gemini_key, prompt_type):
    file_type = file.name.split(".")[-1]  # Get file extension
    
    # Process Video Files
    if file_type in ["mp4", "mov", "avi", "mkv"]:
        st.write("Extracting audio from video...")
        audio_file = extract_audio_from_video(file)  # Extract audio
        
        st.write("Transcribing the audio...")
        transcript = transcribe_audio(assemblyai_key, audio_file)  # Transcribe audio

    # Process Text Files
    elif file_type == "txt":
        transcript = file.read().decode('utf-8')

    # Process Document Files (doc/docx)
    elif file_type in ["doc", "docx"]:
        doc = Document(file)
        transcript = "\n".join([p.text for p in doc.paragraphs])

    else:
        st.error("Unsupported file format.")
        return

    # Generate the appropriate prompt based on user selection
    if prompt_type == "Meeting Notes":
        prompt = generate_insights_prompt()
    elif prompt_type == "Requirement Document":
        prompt = generate_doc_prompt()
    elif prompt_type == "Presentation":
        prompt = generate_presentation_prompt()

    # Ask Google Gemini to generate content based on the prompt and transcript
    st.write("Generating Insights...")
    response = ask_gemini(gemini_key, transcript, prompt)
    
    if response:
        st.write(f"Generated {prompt_type.capitalize()}:\n")
        
        if prompt_type == "Presentation":
            handle_presentation_response(response)  # Generate PowerPoint
        elif prompt_type == "Requirement Document":
            handle_word_doc_response(response)  # Generate Word document
        else:
            handle_text_response(response)  # Generate Text file

        return response
    else:
        st.error(f"Failed to generate {prompt_type.capitalize()}.")

def handle_presentation_response(response):
    response_json = json.loads(response)  # Parse JSON response
    prs = Presentation()

    # Function to add a slide
    def add_slide(prs, title, content):
        slide_layout = prs.slide_layouts[1]  # Use the title + content layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.shapes.placeholders[1].text = content
    
    # Add slides based on the response
    for slide_content in response_json:
        add_slide(prs, slide_content["title"], slide_content["text"])

    # Save the presentation to a BytesIO buffer
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)  # Reset buffer position to the start

    # Create a downloadable link for the presentation with a unique key
    st.download_button(
        label="Download Presentation",
        data=buffer,
        file_name="generated_presentation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        key="download_presentation_button"  # Unique key for this button
    )

def handle_word_doc_response(response):
    cleaned_text = response.replace("```json\n", "").replace("```", "").strip()

    # Try to load the response as JSON. If it's not JSON, treat it as plain text.
    try:
        content = json.loads(cleaned_text)  # Attempt to parse as JSON
    except json.JSONDecodeError:
        st.write("Response is not in JSON format. Treating as plain text.")
        content = {"title": "Generated Document", "sections": [{"heading": "", "body": cleaned_text}]}

    # Create a Word document in memory
    doc = Document()
    
    # Add title to the document
    title = content.get("title", "Document Title")
    title_paragraph = doc.add_heading(level=1)
    title_paragraph.add_run(title).font.size = Pt(24)
    title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Add sections and content to the document
    for section in content.get("sections", []):
        if section.get("heading", ""):
            doc.add_heading(section.get("heading", ""), level=2)
        if section.get("body", ""):
            doc.add_paragraph(section.get("body", ""))

    # Save the document to a BytesIO buffer
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)  # Reset buffer position to the start

    # Create a unique key for the download button using a timestamp
    unique_key = f"download_doc_button_{int(time.time())}"

    # Create a downloadable link for the document with a unique key
    st.download_button(
        label="Download Document",
        data=buffer,
        file_name="generated_document.docx",
        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        key=unique_key  # Unique key for this button
    )

# Handle text response
def handle_text_response(response):
    # Clean and prepare the response text
    cleaned_text = response.strip()

    # Create a downloadable text file
    buffer = BytesIO()
    buffer.write(cleaned_text.encode('utf-8'))  # Write the response as bytes
    buffer.seek(0)  # Reset buffer position to the start

    # Create a unique key for the download button using a timestamp
    unique_key = f"download_meeting_minutes_{int(time.time())}"

    # Create a download button for the text file
    st.download_button(
        label="Download Meeting Minutes",
        data=buffer,
        file_name="generated_output.txt",
        mime="text/plain",
        key=unique_key  # Unique key for this button
    )
# ------------------ Streamlit Interface ------------------

st.title("Meeting Insights Generator")

# API Key Inputs
assemblyai_key = st.text_input("Enter AssemblyAI API Key", type="password")
gemini_key = st.text_input("Enter Google Gemini API Key", type="password")

# Save API keys in session state
if st.button("Save API Keys"):
    save_api_keys(assemblyai_key, gemini_key)
    st.success("API Keys saved!")

# File Upload
uploaded_file = st.file_uploader("Upload Video, Text, or Doc File", type=["mp4", "mov", "avi", "mkv", "txt", "doc", "docx"])

# Select output type (Insights, Document, Presentation)
prompt_type = st.selectbox("Select output type", ["Meeting Notes", "Requirement Document", "Presentation"])

# Process file and generate output
if uploaded_file and st.button("Process File"):
    assemblyai_key, gemini_key = load_api_keys()

    if not assemblyai_key or not gemini_key:
        st.error("API keys are required to proceed.")
    else:
        response = process_file(uploaded_file, assemblyai_key, gemini_key, prompt_type)
        
        # if response:
        #     # Handle Presentation or Document generation
        #     if prompt_type == "Presentation":
        #         handle_presentation_response(response)
        #     elif prompt_type == "Requirement Document":
        #         handle_word_doc_response(response)
